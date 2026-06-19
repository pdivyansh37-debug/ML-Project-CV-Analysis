from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_premium_tonycv_ppt():
    prs = Presentation()
    prs.slide_width = Inches(13.33)  # Widescreen 16:9
    prs.slide_height = Inches(7.5)

    # Palette
    NAVY = RGBColor(15, 23, 42)
    VIOLET = RGBColor(124, 58, 237)
    CYAN = RGBColor(6, 182, 212)
    WHITE = RGBColor(255, 255, 255)
    GRAY = RGBColor(148, 163, 184)
    CARD_BG = RGBColor(30, 41, 59)

    def apply_background(slide):
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = NAVY

    def add_sidebar_accent(slide):
        # Sidebar decorative bar
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.12), Inches(7.5))
        shape.fill.solid()
        shape.fill.fore_color.rgb = VIOLET
        shape.line.fill.background()

    def add_slide_header(slide, title_text, slide_num):
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(10), Inches(1))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.bold = True
        p.font.size = Pt(36)
        p.font.color.rgb = WHITE
        
        # Subtitle badge under title
        badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.4), Inches(2.2), Inches(0.4))
        badge.fill.solid()
        badge.fill.fore_color.rgb = VIOLET
        badge.line.color.rgb = VIOLET
        b_tf = badge.text_frame
        bp = b_tf.paragraphs[0]
        bp.text = "CORE MODULE v2.5"
        bp.font.size = Pt(10)
        bp.font.bold = True
        bp.font.color.rgb = WHITE
        bp.alignment = PP_ALIGN.CENTER

        # Slide Number
        num_box = slide.shapes.add_textbox(Inches(12.5), Inches(0.5), Inches(1), Inches(1))
        np = num_box.text_frame.paragraphs[0]
        np.text = f"{slide_num:02d}"
        np.font.size = Pt(28)
        np.font.bold = True
        np.font.color.rgb = VIOLET
        np.alignment = PP_ALIGN.CENTER

    def add_content_card(slide, top, height, text_content):
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), top, Inches(12), height)
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = RGBColor(51, 65, 85)
        
        tf = card.text_frame
        tf.word_wrap = True
        tf.auto_size = False
        tf.margin_left = Inches(0.3)
        tf.margin_top = Inches(0.2)
        
        for i, line in enumerate(text_content.split('\n')):
            p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
            if line.startswith('•'):
                p.text = line
                p.level = 0
                p.font.size = Pt(20)
                p.font.color.rgb = WHITE
            else:
                p.text = line
                p.font.size = Pt(22)
                p.font.bold = True
                p.font.color.rgb = CYAN
            p.space_after = Pt(12)

    # 1. Cover
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_background(slide)
    
    # Large Center UI Box
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3), Inches(2), Inches(7.33), Inches(3.5))
    box.fill.solid()
    box.fill.fore_color.rgb = CARD_BG
    box.line.color.rgb = VIOLET
    box.line.width = Pt(3)

    title_box = slide.shapes.add_textbox(Inches(3.5), Inches(2.5), Inches(6.33), Inches(2))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "TONYCV"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(80)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    p2 = tf.add_paragraph()
    p2.text = "AI-Powered Recruitment Intelligence"
    p2.alignment = PP_ALIGN.CENTER
    p2.font.size = Pt(24)
    p2.font.color.rgb = CYAN

    author_box = slide.shapes.add_textbox(Inches(3.5), Inches(6), Inches(6.33), Inches(1))
    ap = author_box.text_frame.paragraphs[0]
    ap.text = "Developed by Prashant Singh Rawat"
    ap.alignment = PP_ALIGN.CENTER
    ap.font.size = Pt(14)
    ap.font.color.rgb = GRAY

    # Slides Data
    content_data = [
        ("EXECUTIVE SUMMARY", "PROJECT OVERVIEW\n• End-to-end recruitment platform using NLP/ML.\n• Transforms static CVs into dynamic career roadmaps.\n• Purpose: Bridge the gap between skillsets and industry hiring."),
        ("PROBLEM STATEMENT", "CURRENT CHALLENGES\n• ATS Bias: 75% of resumes are discarded by keywords.\n• Identity Fraud: No verification of technical proficiency.\n• Feedback Gap: Candidates receive rejection without guidance."),
        ("MODELS & TECH", "THE INTELLIGENCE STACK\n• NLP: spaCy (NER) & PDFPlumber for vector space extraction.\n• ML Core: RandomForestClassifier (100 Decision Trees).\n• Architecture: FastAPI (Backend) & React 19 (Frontend)."),
        ("ML WORKFLOW", "THE PIPELINE\n• Input: PDF Document Vectorization.\n• Processing: Entity Recognition & Feature Engineering.\n• Output: Placement Probability & Skill Gap Analysis."),
        ("RESULTS & ACCURACY", "KEY PERFORMANCE DATA\n• 92% Inference Accuracy on synthetic benchmark data.\n• 80% Reduction in manual identity verification time.\n• High precision in identifying missing technical proficiencies."),
        ("FUTURE VISION", "ROADMAP 2026\n• Live Social/GitHub Trace Tracking.\n• LLM-Powered Personalized Interview Coaching.\n• Web3 Blockchain Skill-Verification Minting."),
        ("RESEARCH STATUS", "PATENT & PUBLICATIONS\n• Patent Pending: Multi-Vector Verification Identity Logic.\n• Targeted: Publication in top-tier AI/IEEE Journals.\n• Status: Ready for Enterprise Scaling.")
    ]

    for i, (title, content) in enumerate(content_data):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        apply_background(slide)
        add_sidebar_accent(slide)
        add_slide_header(slide, title, i + 2)
        add_content_card(slide, Inches(2.2), Inches(4.5), content)

    prs.save('TonyCV_Attractive_Presentation.pptx')
    print("Premium Presentation Generated.")

if __name__ == "__main__":
    create_premium_tonycv_ppt()
